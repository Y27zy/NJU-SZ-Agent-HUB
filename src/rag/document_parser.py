from pathlib import Path


def parse_document(file_path: str | Path) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        import fitz

        text_parts: list[str] = []
        with fitz.open(path) as doc:
            for page in doc:
                text_parts.append(page.get_text())
        return "\n".join(text_parts)
    if suffix == ".pptx":
        from pptx import Presentation

        prs = Presentation(path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    if suffix in {".txt", ".md", ".markdown"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"Unsupported file type: {suffix}")
